from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

NAVY = RGBColor(0x2A, 0x31, 0x42)
GOLD = RGBColor(0xC5, 0xA5, 0x5A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
RED = RGBColor(0xDC, 0x45, 0x45)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
LIGHT_GRAY = RGBColor(0xE8, 0xE3, 0xDA)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_footer(slide, text="BITE ME · PRICING SIMULATION", page_num=None):
    add_text(slide, 0.6, 6.9, 4, 0.4, text, size=8, color=GRAY)
    if page_num:
        add_text(slide, 11.5, 0.2, 1.5, 0.3, f"{page_num} / 10", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def add_section_indicator(slide, label):
    add_rect(slide, 0.6, 0.5, 0.06, 0.3, GOLD)
    add_text(slide, 0.8, 0.48, 5, 0.35, label, size=9, color=GRAY, font_name="Consolas")

def add_kpi_card(slide, left, top, width, label, value, sub="", border_color=GOLD, value_color=NAVY):
    card = add_rect(slide, left, top, width, 0.9, WHITE)
    bar = add_rect(slide, left, top, 0.04, 0.9, border_color)
    add_text(slide, left + 0.15, top + 0.05, width - 0.2, 0.2, label, size=8, color=GRAY)
    add_text(slide, left + 0.15, top + 0.28, width - 0.2, 0.35, value, size=18, color=value_color, bold=True)
    if sub:
        add_text(slide, left + 0.15, top + 0.65, width - 0.2, 0.2, sub, size=8, color=GRAY)

def add_table(slide, left, top, width, headers, rows, col_widths=None):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(0.3 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Malgun Gothic"
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            text = str(val).replace("✓", "O").replace("×", "X")
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = NAVY
            p.font.name = "Malgun Gothic"
            p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if "-" in text and ci > 0 and any(c.isdigit() for c in text):
                p.font.color.rgb = RED
            elif "+" in text and ci > 0 and any(c.isdigit() for c in text):
                p.font.color.rgb = GREEN

    return table_shape

def add_takeaway(slide, top, title, body_lines):
    add_rect(slide, 0.6, top, 12.1, 0.04, GOLD)
    add_rect(slide, 0.6, top + 0.04, 12.1, len(body_lines) * 0.28 + 0.6, NAVY)
    add_text(slide, 0.9, top + 0.12, 5, 0.25, title, size=9, color=GOLD, bold=True, font_name="Consolas")
    y = top + 0.38
    for line in body_lines:
        add_text(slide, 0.9, y, 11.5, 0.25, line, size=11, color=WHITE)
        y += 0.26


# ========== PAGE 1: COVER ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, GOLD)
add_text(slide, 1.0, 1.5, 8, 0.4, "P R I C I N G   S I M U L A T I O N", size=11, color=GOLD, font_name="Consolas")
add_text(slide, 1.0, 2.2, 10, 0.8, "보솜 라이트", size=36, color=WHITE, bold=True)
add_text(slide, 1.0, 3.0, 10, 0.7, "원가 인상 대응 PQ 시뮬레이션", size=30, color=WHITE, bold=True)
add_text(slide, 1.0, 3.8, 10, 0.5, "가격-물량 시나리오 분석 및 세트 플레이 전략 권장안", size=14, color=RGBColor(0x99, 0x99, 0xAA))
add_text(slide, 1.0, 5.5, 6, 0.3, "2026.05 | BITE ME Brand Team", size=11, color=RGBColor(0x66, 0x66, 0x77))


# ========== PAGE 2: EXECUTIVE SUMMARY ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "E X E C U T I V E   S U M M A R Y")
add_text(slide, 0.6, 0.9, 10, 0.5, "핵심 요약", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=2)

add_takeaway(slide, 1.5, "KEY TAKEAWAY", [
    "보솜 라이트는 원가 인상(+6.3~6.9%) 이전부터 구조적 적자 상태이며,",
    "S3(+2,000원) 전면 인상 + 세트 플레이 전략을 권장합니다.",
    "100% 세트 전환 시 손익분기 할인율 14.5% — 변동비 30% 이하 관리가 핵심."
])

add_kpi_card(slide, 0.6, 3.1, 2.8, "현재 단위당 CM", "-271원", "원가 인상 전 기준", RED, RED)
add_kpi_card(slide, 3.6, 3.1, 2.8, "원가 인상 후 CM", "-650원", "가격 미조정 시", RED, RED)
add_kpi_card(slide, 6.6, 3.1, 2.8, "손익분기 최소 인상", "+1,400원", "변동비 30% 기준", GOLD, GOLD)
add_kpi_card(slide, 9.6, 3.1, 2.8, "S3 적용 시 월간 CM", "+223만원", "+2,000원 / 물량 85%", GREEN, GREEN)

add_text(slide, 0.6, 4.4, 5, 0.35, "권장 액션", size=14, color=NAVY, bold=True)

actions = [
    "· 단품 +2,000원 인상 (대형 14,500 / 표준형 15,500)",
    "· 세트 플레이: 단품 앵커가로 4팩 전환 유도",
    "· 세트 할인율 14% 이하 유지 (손익분기 14.5%)",
]
for i, a in enumerate(actions):
    add_text(slide, 0.8, 4.8 + i * 0.3, 5.5, 0.3, a, size=11, color=NAVY)

actions2 = [
    "· 인상 후 2~4주 물량 모니터링",
    "· 변동비 30% 이하 관리 (세트 플레이 성패 좌우)",
    "· 중장기: 대나무/재생펄프 패드로 라이트 교체 검토",
]
for i, a in enumerate(actions2):
    add_text(slide, 6.8, 4.8 + i * 0.3, 5.5, 0.3, a, size=11, color=NAVY)


# ========== PAGE 3: BACKGROUND ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "B A C K G R O U N D")
add_text(slide, 0.6, 0.9, 10, 0.5, "배경: 원가 인상 & 현재 손익 구조", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=3)

add_text(slide, 0.6, 1.5, 5, 0.3, "SKU별 원가 인상 현황", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 1.9, 12.1,
    ["구분", "판매가", "현재 원가", "5월 예상", "인상분", "인상율", "매출총이익(현재)", "매출총이익(인상후)"],
    [
        ["대형 단품", "12,500", "5,523", "5,901", "+378", "6.9%", "6,977", "6,599"],
        ["대형 4팩", "42,900", "22,094", "23,606", "+1,512", "6.9%", "20,806", "19,294"],
        ["표준형 단품", "13,500", "6,061", "6,441", "+380", "6.3%", "7,439", "7,059"],
        ["표준형 4팩", "45,900", "24,246", "25,766", "+1,520", "6.3%", "21,654", "20,134"],
    ])

add_text(slide, 0.6, 4.0, 5, 0.3, "2026 Jan-Apr 월별 실적", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 4.35, 12.1,
    ["월", "출고량", "순매출", "매출원가", "매출총이익", "공헌이익"],
    [
        ["1월", "5,284", "5,330만", "3,593만", "1,736만", "-1.5만"],
        ["2월", "12,922", "7,519만", "5,399만", "2,120만", "-273만"],
        ["3월", "9,789", "9,190만", "6,495만", "2,695만", "-57만"],
        ["4월", "3,852", "3,444만", "3,211만", "233만", "-18만"],
        ["월평균", "7,962", "6,371만", "4,675만", "1,696만", "-87만"],
    ])


# ========== PAGE 4: ASSUMPTIONS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "A S S U M P T I O N S")
add_text(slide, 0.6, 0.9, 10, 0.5, "시뮬레이션 핵심 가정", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=4)

add_kpi_card(slide, 0.6, 1.5, 3.6, "변동비율 (순매출 대비)", "30%", "물류/수수료/프로모션 등")
add_kpi_card(slide, 4.5, 1.5, 3.6, "순매출 전환율 (판매가 대비)", "70%", "채널 믹스 가중평균")
add_kpi_card(slide, 8.4, 1.5, 3.6, "월평균 출고량", "7,962", "대형 58% / 표준형 42%")

add_text(slide, 0.6, 2.8, 5, 0.3, "단위 경제 구조 (개당)", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 3.15, 8,
    ["항목", "금액", "비고"],
    [
        ["순매출", "8,001원", "판매가 대비 ~70%"],
        ["변동비 (30%)", "-2,400원", "물류/수수료/프로모션"],
        ["매출원가 (인상 후)", "-6,251원", "기존 5,872 + 인상분 379"],
        ["공헌이익 (CM)", "-650원", "가격 미조정 시"],
    ])

add_text(slide, 0.6, 5.1, 5, 0.3, "순매출 전환율 70% 산출 근거", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 5.45, 8,
    ["채널", "추정 전환율", "순매출 증분", "비고"],
    [
        ["자사몰", "~85-90%", "+1,700~1,800원", "할인/쿠폰/결제수수료"],
        ["쿠팡 (공급가)", "~60%", "+1,200원", "공급가 기준 출고"],
        ["가중평균", "~70%", "~+1,400원", "채널 믹스 반영"],
    ])


# ========== PAGE 5: SCENARIOS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "S C E N A R I O   A N A L Y S I S")
add_text(slide, 0.6, 0.9, 10, 0.5, "3대 시나리오 비교", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=5)

scenarios = [
    ("S1", "소폭 인상", "+1,000원", "-160원", "월 -108만", "적자 축소", LIGHT_GRAY, RED),
    ("S2", "손익분기", "+1,500원", "+85원", "월 +58만", "손익분기 돌파", LIGHT_GRAY, GREEN),
    ("S3 BEST", "최대 인상", "+2,000원", "+330원", "월 +223만", "흑자 전환", GOLD, GREEN),
]
for i, (tag, name, price, cm, monthly, verdict, border, v_color) in enumerate(scenarios):
    x = 0.6 + i * 4.1
    card = add_rect(slide, x, 1.5, 3.8, 2.3, WHITE, border)
    add_text(slide, x + 0.2, 1.6, 1.5, 0.25, tag, size=9, color=GOLD, bold=True, font_name="Consolas")
    add_text(slide, x + 0.2, 1.85, 3, 0.25, name, size=12, color=GRAY)
    add_text(slide, x + 0.2, 2.15, 3, 0.4, price, size=26, color=NAVY, bold=True)
    add_text(slide, x + 0.2, 2.6, 3, 0.2, f"단위 CM: {cm}", size=10, color=GRAY)
    add_text(slide, x + 0.2, 2.85, 3, 0.2, f"85% 유지 시", size=10, color=GRAY)
    add_text(slide, x + 0.2, 3.1, 3, 0.3, monthly, size=15, color=v_color, bold=True)
    add_text(slide, x + 0.2, 3.45, 3, 0.2, verdict, size=10, color=v_color, bold=True)

add_text(slide, 0.6, 4.1, 5, 0.3, "시나리오별 세부 비교 (물량 유지율 85%)", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 4.45, 12.1,
    ["항목", "현재(인상전)", "가격미조정", "S1 +1,000", "S2 +1,500", "S3 +2,000"],
    [
        ["순매출/개", "8,001", "8,001", "8,701", "9,051", "9,401"],
        ["변동비/개(30%)", "2,400", "2,400", "2,610", "2,715", "2,820"],
        ["원가/개", "5,872", "6,251", "6,251", "6,251", "6,251"],
        ["CM/개", "-271", "-650", "-160", "+85", "+330"],
        ["월간 CM", "-216만", "-518만", "-108만", "+58만", "+223만"],
        ["연간 CM", "-2,593만", "-6,210만", "-1,296만", "+690만", "+2,680만"],
    ])


# ========== PAGE 6: SENSITIVITY ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "S E N S I T I V I T Y")
add_text(slide, 0.6, 0.9, 10, 0.5, "민감도 분석", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=6)

add_text(slide, 0.6, 1.5, 10, 0.3, "가격 인상 x 물량 유지율 → 월간 공헌이익 (만원)", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 1.9, 12.1,
    ["인상분\\유지율", "100%", "95%", "90%", "85%", "80%", "75%", "70%", "60%", "50%"],
    [
        ["미조정", "-518", "-492", "-466", "-440", "-414", "-388", "-363", "-311", "-259"],
        ["+500원", "-450", "-428", "-405", "-383", "-360", "-338", "-315", "-270", "-225"],
        ["+1,000원(S1)", "-188", "-179", "-169", "-160", "-150", "-141", "-132", "-113", "-94"],
        ["+1,200원", "-116", "-110", "-104", "-99", "-93", "-87", "-81", "-70", "-58"],
        ["+1,400원(BEP)", "+3", "+3", "+2", "+2", "+2", "+2", "+2", "+1", "+1"],
        ["+1,500원(S2)", "+68", "+64", "+61", "+58", "+54", "+51", "+47", "+41", "+34"],
        ["+1,800원", "+176", "+167", "+158", "+150", "+141", "+132", "+123", "+106", "+88"],
        ["+2,000원(S3)", "+263", "+250", "+237", "+223", "+210", "+197", "+184", "+158", "+131"],
    ])

add_text(slide, 0.6, 4.2, 10, 0.3, "S3(+2,000원) 물량 시나리오별 연간 CM", size=13, color=NAVY, bold=True)
add_kpi_card(slide, 0.6, 4.6, 2.8, "100% 유지", "+3,153만", "월 263만", GREEN, GREEN)
add_kpi_card(slide, 3.6, 4.6, 2.8, "85% 유지", "+2,680만", "월 223만", GREEN, GREEN)
add_kpi_card(slide, 6.6, 4.6, 2.8, "70% 유지", "+2,207만", "월 184만", GREEN, GREEN)
add_kpi_card(slide, 9.6, 4.6, 2.8, "50% 유지", "+1,576만", "월 131만", GREEN, GREEN)


# ========== PAGE 7: SKU PRICING ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "P R I C I N G   D E T A I L")
add_text(slide, 0.6, 0.9, 10, 0.5, "SKU별 가격 설정 & 세트 플레이 전략", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=7)

add_text(slide, 0.6, 1.5, 5, 0.3, "S3(+2,000원) 적용 시 가격표", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 1.9, 12.1,
    ["구분", "현재 판매가", "S3 판매가", "인상분", "타사 최고가", "판단"],
    [
        ["대형 단품", "12,500", "14,500", "+2,000", "14,500~14,900", "안전"],
        ["대형 4팩", "42,900", "49,900", "+7,000", "--", "--"],
        ["표준형 단품", "13,500", "15,500", "+2,000", "15,500~15,900", "안전"],
        ["표준형 4팩", "45,900", "53,300", "+7,400", "--", "--"],
    ])

add_takeaway(slide, 3.6, "SET PLAY STRATEGY", [
    "단품 +2,000원을 앵커가로 설정 → 4팩 세트의 가격 메리트 부각 → 100% 세트 판매 구조로 자연 전환",
    "단품 폐지를 직접 실행하지 않고, 가격 차등으로 고객이 스스로 세트를 선택하도록 유도",
    "",
    "핵심 질문: 100% 4팩 판매 전환 시, 할인율을 몇 %까지 줄 수 있으면서 흑자를 유지하는가?",
    "→ 손익분기 할인율 14.5% / 상세 분석은 다음 페이지 참조",
])


# ========== PAGE 8: SET PLAY ANALYSIS ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "S E T   P L A Y   A N A L Y S I S")
add_text(slide, 0.6, 0.9, 12, 0.5, "세트 플레이 전략: 할인율별 손익 분석", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=8)

add_text(slide, 0.6, 1.45, 12, 0.25, "전제: 단품 +2,000원 앵커 후 100% 4팩 전환 | 가중평균 단품가 14,915원 | 변동비 30% | 전환율 70%", size=9, color=GRAY)

add_text(slide, 0.6, 1.8, 5, 0.3, "세트 할인율별 단위 경제성", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 2.15, 12.1,
    ["할인율", "4팩 단가(개당)", "순매출/개", "변동비", "원가", "CM/개", "월간 CM(85%)", "판정"],
    [
        ["5%", "14,169", "9,918", "2,975", "6,251", "+692", "+468만", "안정 흑자"],
        ["8%", "13,722", "9,605", "2,882", "6,251", "+473", "+320만", "안정 흑자"],
        ["10%", "13,424", "9,397", "2,819", "6,251", "+327", "+221만", "안정 흑자"],
        ["12%", "13,125", "9,188", "2,756", "6,251", "+180", "+122만", "흑자"],
        ["14% ◀현재", "12,827", "8,979", "2,694", "6,251", "+34", "+23만", "손익분기 근접"],
        ["15%", "12,678", "8,875", "2,662", "6,251", "-39", "-26만", "적자"],
        ["18%", "12,230", "8,561", "2,568", "6,251", "-258", "-175만", "적자"],
        ["20%", "11,932", "8,352", "2,506", "6,251", "-404", "-273만", "적자"],
    ])

add_kpi_card(slide, 0.6, 5.1, 3.6, "손익분기 할인율", "14.5%", "이 이하에서 흑자", GOLD, GOLD)
add_kpi_card(slide, 4.5, 5.1, 3.6, "14% 할인 시 CM/개", "+34원", "현재 할인율 유지 기준", GREEN, GREEN)
add_kpi_card(slide, 8.4, 5.1, 3.6, "안전 마진 할인율", "12%", "CM >= 200원/개", GREEN, GREEN)

add_takeaway(slide, 6.15, "KEY INSIGHT", [
    "14% 할인 유지 시 흑자(+34원)이나 마진 타이트. 할인율 12% 이하 시 안정적. 변동비 관리가 핵심.",
])


# ========== PAGE 9: SET PLAY BEP BY PRICE ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "S E T   P L A Y   B E P   B Y   P R I C E")
add_text(slide, 0.6, 0.9, 12, 0.5, "인상분별 세트 BEP 할인율 비교", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=9)

add_text(slide, 0.6, 1.45, 12, 0.25, "인상폭이 낮으면 앵커가가 낮아져 세트 BEP 할인율이 급격히 하락 — +2,000원(S3)에서만 14% 할인 유지 가능", size=9, color=GRAY)

add_text(slide, 0.6, 1.8, 5, 0.3, "인상분 → 세트 할인율 손익분기점", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 2.15, 12.1,
    ["인상분", "가중평균 단품가", "세트 BEP 할인율", "14% 할인 시 CM/개", "14% 할인 월간CM", "판정"],
    [
        ["+1,000원 (S1)", "13,915원", "8.3%", "-387원", "-262만", "14% 할인 불가"],
        ["+1,200원", "14,115원", "9.6%", "-303원", "-205만", "14% 할인 불가"],
        ["+1,400원 (BEP)", "14,315원", "10.9%", "-219원", "-148만", "14% 할인 불가"],
        ["+1,500원 (S2)", "14,415원", "11.5%", "-176원", "-119만", "14% 할인 불가"],
        ["+1,800원", "14,715원", "13.3%", "-50원", "-34만", "14% 할인 근접"],
        ["+2,000원 (S3)", "14,915원", "14.5%", "+34원", "+23만", "14% 할인 가능"],
    ])

add_kpi_card(slide, 0.6, 4.6, 3.6, "S2(+1,500) BEP 할인율", "11.5%", "14% 할인 시 적자 -176원/개", RED, RED)
add_kpi_card(slide, 4.5, 4.6, 3.6, "S3(+2,000) BEP 할인율", "14.5%", "14% 할인 시 흑자 +34원/개", GOLD, GOLD)
add_kpi_card(slide, 8.4, 4.6, 3.6, "S3 + 할인율 12%", "+180원/개", "안정적 마진 확보", GREEN, GREEN)

add_takeaway(slide, 5.8, "KEY INSIGHT", [
    "+2,000원(S3)에서만 14% 할인 가능. S2(+1,500)는 할인율 11% 이하 필요 → 세트 메리트 약화.",
    "앵커가가 높아야 할인율 여유가 생기고, 세트 플레이가 성립합니다.",
])


# ========== PAGE 10: RECOMMENDATION ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_section_indicator(slide, "R E C O M M E N D A T I O N")
add_text(slide, 0.6, 0.9, 10, 0.5, "권장안 & Next Steps", size=22, color=NAVY, bold=True)
add_footer(slide, page_num=10)

add_takeaway(slide, 1.5, "RECOMMENDATION", [
    "S3: 단품 +2,000원 앵커가 + 세트 플레이 전략",
    "대형 14,500 / 표준형 15,500 (단품 앵커) → 4팩 세트 할인율 14% 이하 유지 → 100% 세트 전환 유도",
])

add_kpi_card(slide, 0.6, 2.7, 2.8, "단품 인상 효과 (85%)", "+223만/월", "현 믹스 기준", GREEN, GREEN)
add_kpi_card(slide, 3.6, 2.7, 2.8, "100% 세트 전환 시", "+23만/월", "14% 할인 유지 기준", GREEN, GREEN)
add_kpi_card(slide, 6.6, 2.7, 2.8, "세트 손익분기 할인율", "14.5%", "이 이하에서 흑자", GOLD, GOLD)
add_kpi_card(slide, 9.6, 2.7, 2.8, "안전 마진 할인율", "12%", "CM >= 200원/개", GOLD, GOLD)

add_text(slide, 0.6, 4.0, 5, 0.3, "실행 로드맵", size=13, color=NAVY, bold=True)
add_table(slide, 0.6, 4.35, 12.1,
    ["단계", "시기", "액션", "판단 기준"],
    [
        ["1. 가격 인상", "즉시", "단품 +2,000원(앵커) / 4팩 할인율 14% 이하 연동", "--"],
        ["2. 모니터링", "인상 후 2~4주", "물량 추적 + 단품/세트 비중 변화 확인", "세트 전환율 상승 여부"],
        ["3. 세트 전환 판단", "인상 후 1개월", "세트 비중 80%+ → 세트 플레이 정착", "월간 CM 흑자 + 세트 비중"],
        ["4. 할인율 최적화", "인상 후 2~3개월", "세트 전환 정착 시 할인율 12%로 하향 검토", "CM >= 200원/개 목표"],
        ["5. 중장기", "Q3~Q4", "대나무/재생펄프 패드 기획, 라이트 대체", "원가 대비 시장가 경쟁력"],
    ])


# Save
output_path = r"C:\Users\User\Downloads\보솜라이트_PQ시뮬레이션_보고서.pptx"
prs.save(output_path)
import os
print(f"PPTX saved: {output_path}")
print(f"Size: {os.path.getsize(output_path):,} bytes")
print(f"Slides: {len(prs.slides)}")
